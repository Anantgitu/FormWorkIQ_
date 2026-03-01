"""
FormworkIQ — L&T CreaTech 2026 PPT Generator
Generates a professional 8-slide PPTX presentation using python-pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── L&T Brand Colors ──────────────────────────────────────────────
LT_RED    = RGBColor(0xC8, 0x10, 0x2E)   # L&T primary red
LT_NAVY   = RGBColor(0x00, 0x33, 0x66)   # L&T dark navy
LT_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LT_GRAY   = RGBColor(0xF0, 0xF0, 0xF0)
LT_DGRAY  = RGBColor(0x33, 0x33, 0x33)
LT_LGRAY  = RGBColor(0x78, 0x78, 0x78)
LT_GREEN  = RGBColor(0x10, 0x80, 0x40)   # for savings highlights

W = Inches(13.33)   # widescreen 16:9 width
H = Inches(7.5)     # widescreen 16:9 height

# ── Helper functions ──────────────────────────────────────────────
def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, text, l, t, w, h,
                 size=18, bold=False, color=LT_WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]  # blank
    return prs.slides.add_slide(layout)

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

# ── Slide builders ────────────────────────────────────────────────

def slide_title(prs):
    slide = add_slide(prs)
    set_bg(slide, LT_NAVY)

    # Red accent bar top
    add_rect(slide, 0, 0, W, Inches(0.08), LT_RED)

    # White diagonal panel (right side decoration)
    add_rect(slide, Inches(8.5), 0, Inches(4.83), H, RGBColor(0x00, 0x28, 0x50))

    # Red vertical stripe
    add_rect(slide, Inches(8.4), 0, Inches(0.12), H, LT_RED)

    # L&T Logo text top-left
    add_text_box(slide, "L&T CreaTech 2026", Inches(0.5), Inches(0.2), Inches(4), Inches(0.5),
                 size=13, bold=True, color=LT_RED)
    add_text_box(slide, "Problem Statement 4", Inches(0.5), Inches(0.65), Inches(5), Inches(0.4),
                 size=11, color=LT_LGRAY)

    # Main title
    add_text_box(slide, "FormworkIQ", Inches(0.5), Inches(1.6), Inches(7.5), Inches(1.4),
                 size=60, bold=True, color=LT_WHITE)
    add_text_box(slide, "AI-Driven Formwork Kitting &\nBoQ Optimization System",
                 Inches(0.5), Inches(3.0), Inches(7.5), Inches(1.2),
                 size=22, bold=False, color=LT_GRAY)

    # Tagline
    add_text_box(slide, "Eliminating over-ordering. Maximizing panel reuse.\nReducing project cost by up to 40%.",
                 Inches(0.5), Inches(4.3), Inches(7.5), Inches(1.0),
                 size=14, italic=True, color=RGBColor(0xB0, 0xC8, 0xE8))

    # Bottom bar
    add_rect(slide, 0, Inches(7.1), W, Inches(0.4), LT_RED)
    add_text_box(slide, "Powered by 2D Bin Packing  ·  Repetition Matching  ·  MERN + FastAPI",
                 Inches(0.3), Inches(7.1), W - Inches(0.6), Inches(0.4),
                 size=10, color=LT_WHITE, align=PP_ALIGN.CENTER)

def slide_problem(prs):
    slide = add_slide(prs)
    set_bg(slide, LT_WHITE)
    add_rect(slide, 0, 0, W, Inches(1.3), LT_NAVY)
    add_text_box(slide, "The Problem", Inches(0.5), Inches(0.2), Inches(10), Inches(0.6),
                 size=30, bold=True, color=LT_WHITE)
    add_text_box(slide, "Why does formwork cost 7–10% of total project cost — and overrun by 30% every time?",
                 Inches(0.5), Inches(0.78), Inches(11), Inches(0.45),
                 size=13, color=LT_GRAY, italic=True)

    problems = [
        ("📋", "Manual BoQ Planning",      "Site engineers use Excel. No algorithm. Massive over-ordering of panels is the norm."),
        ("♻️", "Zero Reuse Tracking",       "Freed panels from completed floors sit idle. Nobody knows which panels can move to next pour."),
        ("📦", "No Kitting System",         "Workers waste hours locating panels on-site. Missing connectors cause pour delays."),
        ("🔍", "No Panel Health Tracking",  "Panels past their reuse limit (50x) stay in rotation. Structural risk + cost hidden."),
    ]

    for i, (icon, title, desc) in enumerate(problems):
        x = Inches(0.4 + (i % 2) * 6.2)
        y = Inches(1.5 + (i // 2) * 2.6)
        add_rect(slide, x, y, Inches(5.8), Inches(2.3), RGBColor(0xF7, 0xF7, 0xF7))
        add_rect(slide, x, y, Inches(0.07), Inches(2.3), LT_RED)
        add_text_box(slide, f"{icon}  {title}", x + Inches(0.15), y + Inches(0.12),
                     Inches(5.5), Inches(0.45), size=15, bold=True, color=LT_NAVY)
        add_text_box(slide, desc, x + Inches(0.15), y + Inches(0.6),
                     Inches(5.5), Inches(1.5), size=12, color=LT_DGRAY)

    add_rect(slide, 0, Inches(7.1), W, Inches(0.4), LT_RED)

def slide_solution(prs):
    slide = add_slide(prs)
    set_bg(slide, LT_WHITE)
    add_rect(slide, 0, 0, W, Inches(1.3), LT_RED)
    add_text_box(slide, "FormworkIQ — Our Solution", Inches(0.5), Inches(0.25), Inches(12), Inches(0.65),
                 size=30, bold=True, color=LT_WHITE)
    add_text_box(slide, "One automated pipeline: from BIM data → optimized BoQ → daily kitting list → panel health loop",
                 Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
                 size=12, color=RGBColor(0xFF, 0xD0, 0xD5), italic=True)

    steps = [
        ("1", "Upload Schedule",  "Engineer uploads building schedule JSON\n(from Revit/Primavera or manual entry)"),
        ("2", "2D Bin Packing",   "Algorithm fills each element's area with\noptimal standard panel combinations"),
        ("3", "Repetition Match", "Identifies freed panels from prior pours\nand schedules reuse — no new order needed"),
        ("4", "Generate BoQ",     "Baseline vs optimized BoQ with exact ₹\nsavings per panel type, ready to procure"),
        ("5", "Daily Kitting",    "Auto-generates per-day worker picking\nlists — panel IDs, locations, connectors"),
        ("6", "QR Feedback",      "Site scans update panel health score.\nAuto-alerts trigger replacement orders"),
    ]

    for i, (num, title, desc) in enumerate(steps):
        x = Inches(0.3 + (i % 3) * 4.3)
        y = Inches(1.55 + (i // 3) * 2.8)
        add_rect(slide, x, y, Inches(4.0), Inches(2.5), RGBColor(0xF5, 0xF8, 0xFF))
        add_rect(slide, x, y, Inches(4.0), Inches(0.5), LT_NAVY)
        add_text_box(slide, f"Step {num}  ·  {title}", x + Inches(0.1), y + Inches(0.04),
                     Inches(3.8), Inches(0.42), size=13, bold=True, color=LT_WHITE)
        add_text_box(slide, desc, x + Inches(0.15), y + Inches(0.6),
                     Inches(3.7), Inches(1.8), size=12, color=LT_DGRAY)

    add_rect(slide, 0, Inches(7.1), W, Inches(0.4), LT_NAVY)

def slide_architecture(prs):
    slide = add_slide(prs)
    set_bg(slide, LT_WHITE)
    add_rect(slide, 0, 0, W, Inches(1.1), LT_NAVY)
    add_text_box(slide, "System Architecture", Inches(0.5), Inches(0.2), Inches(10), Inches(0.6),
                 size=30, bold=True, color=LT_WHITE)
    add_text_box(slide, "Microservices: React Frontend  →  Node.js API  →  Python Engine  →  MongoDB",
                 Inches(0.5), Inches(0.77), Inches(12), Inches(0.32),
                 size=12, color=LT_GRAY, italic=True)

    layers = [
        (LT_RED,  "⚛️  React Frontend (Vite)",      "Port 5173",
         "• 8 pages: Home, Upload, Dashboard\n• Recharts bar charts + KPI cards\n• Framer Motion animations\n• QR scan simulation page"),
        (LT_NAVY, "🟢  Node.js + Express Backend",  "Port 5000",
         "• REST API gateway\n• Proxies to Python optimizer\n• Mongoose models: Inventory, BoQ\n• dotenv config management"),
        (RGBColor(0x8B,0x00,0x00), "🐍  Python FastAPI Engine",  "Port 8000",
         "• 2D Bin Packing algorithm\n• Repetition Matcher engine\n• BoQ generator & kitting\n• Resource leveling logic"),
        (RGBColor(0x00,0x60,0x00), "🍃  MongoDB Atlas",          "Cloud",
         "• OptimizedKit collection\n• LiveInventory tracking\n• FormworkCatalog seeding\n• Flexible nested schemas"),
    ]

    for i, (col, title, port, desc) in enumerate(layers):
        x = Inches(0.3 + i * 3.2)
        y = Inches(1.3)
        add_rect(slide, x, y, Inches(3.0), Inches(5.8), RGBColor(0xF7,0xF7,0xF7))
        add_rect(slide, x, y, Inches(3.0), Inches(0.7), col)
        add_text_box(slide, title, x + Inches(0.1), y + Inches(0.05),
                     Inches(2.8), Inches(0.45), size=12, bold=True, color=LT_WHITE)
        add_text_box(slide, port, x + Inches(0.1), y + Inches(0.5),
                     Inches(2.8), Inches(0.25), size=10, color=RGBColor(0xDD,0xDD,0xDD))
        add_text_box(slide, desc, x + Inches(0.1), y + Inches(0.9),
                     Inches(2.8), Inches(4.5), size=11, color=LT_DGRAY)

        # Arrow between boxes
        if i < 3:
            add_text_box(slide, "→", x + Inches(3.0), y + Inches(2.5),
                         Inches(0.2), Inches(0.5), size=20, bold=True, color=LT_RED,
                         align=PP_ALIGN.CENTER)

    add_rect(slide, 0, Inches(7.1), W, Inches(0.4), LT_RED)

def slide_algorithm(prs):
    slide = add_slide(prs)
    set_bg(slide, LT_WHITE)
    add_rect(slide, 0, 0, W, Inches(1.1), LT_RED)
    add_text_box(slide, "Core Algorithm: How the Optimizer Works", Inches(0.5), Inches(0.2),
                 Inches(12), Inches(0.65), size=28, bold=True, color=LT_WHITE)

    # Left: 2D Bin Packing
    add_rect(slide, Inches(0.3), Inches(1.3), Inches(6.0), Inches(5.8), RGBColor(0xF0,0xF4,0xFF))
    add_rect(slide, Inches(0.3), Inches(1.3), Inches(6.0), Inches(0.55), LT_NAVY)
    add_text_box(slide, "🧠  2D Bin Packing", Inches(0.4), Inches(1.34),
                 Inches(5.8), Inches(0.45), size=16, bold=True, color=LT_WHITE)

    bp_text = (
        "For each structural element (Column, Wall, Core-Wall):\n\n"
        "1. Calculate required formwork area = width × height\n\n"
        "2. Greedily fill using standard panels:\n"
        "   P-2×1 (2m²) → P-1×1 (1m²) → P-0.5×1 → P-0.3×1\n\n"
        "3. Use filler strips for remainder areas\n\n"
        "4. Record exact panel count needed per element\n\n"
        "📐 Example: 5m × 3m Wall = 15m²\n"
        "   → 7× P-2×1 + 1× P-1×1 = 15m² covered\n"
        "   → Optimal combination, min number of pieces"
    )
    add_text_box(slide, bp_text, Inches(0.45), Inches(2.0),
                 Inches(5.7), Inches(5.0), size=11.5, color=LT_DGRAY)

    # Right: Repetition Matcher
    add_rect(slide, Inches(6.9), Inches(1.3), Inches(6.0), Inches(5.8), RGBColor(0xFF,0xF4,0xF0))
    add_rect(slide, Inches(6.9), Inches(1.3), Inches(6.0), Inches(0.55), LT_RED)
    add_text_box(slide, "♻️  Repetition Matcher", Inches(7.0), Inches(1.34),
                 Inches(5.8), Inches(0.45), size=16, bold=True, color=LT_WHITE)

    rm_text = (
        "Eliminates redundant material purchases:\n\n"
        "1. Sort all pours chronologically by pour_date\n\n"
        "2. For each element, compute available_date:\n"
        "   pour_date + cure_days + 1 cleaning day\n\n"
        "3. Before ordering new panels, check if:\n"
        "   freed panels ≥ required panels (same type)\n\n"
        "4. Tag kit as REUSE, PARTIAL_REUSE, or NEW_ORDER\n\n"
        "💰 Impact Example:\n"
        "   Floor 1 Wall freed → Floor 2 Wall reuses same\n"
        "   panels → 0 new panels ordered → 100% saving\n"
        "   for that element"
    )
    add_text_box(slide, rm_text, Inches(7.0), Inches(2.0),
                 Inches(5.7), Inches(5.0), size=11.5, color=LT_DGRAY)

    add_rect(slide, 0, Inches(7.1), W, Inches(0.4), LT_NAVY)

def slide_results(prs):
    slide = add_slide(prs)
    set_bg(slide, LT_WHITE)
    add_rect(slide, 0, 0, W, Inches(1.1), LT_NAVY)
    add_text_box(slide, "Results & Impact", Inches(0.5), Inches(0.2),
                 Inches(10), Inches(0.65), size=30, bold=True, color=LT_WHITE)
    add_text_box(slide, "Simulated on 550-element, 10-floor mega-structure (LNT-TOWER-OMEGA)",
                 Inches(0.5), Inches(0.77), Inches(12), Inches(0.32),
                 size=12, color=LT_GRAY, italic=True)

    kpis = [
        ("₹4.2L+",  "Cost Savings",          LT_GREEN),
        ("35–40%",  "BoQ Reduction",          LT_RED),
        ("180+",    "Panels Eliminated",      LT_NAVY),
        ("100%",    "Daily Kits Automated",   RGBColor(0x80,0x40,0x00)),
    ]
    for i, (val, label, col) in enumerate(kpis):
        x = Inches(0.3 + i * 3.2)
        add_rect(slide, x, Inches(1.25), Inches(3.0), Inches(1.9), RGBColor(0xF5,0xF5,0xF5))
        add_rect(slide, x, Inches(1.25), Inches(3.0), Inches(0.07), col)
        add_text_box(slide, val, x, Inches(1.35), Inches(3.0), Inches(1.0),
                     size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text_box(slide, label, x, Inches(2.28), Inches(3.0), Inches(0.4),
                     size=13, color=LT_LGRAY, align=PP_ALIGN.CENTER)

    # Feature highlights
    features = [
        ("🧠", "Algorithm", "2D Bin Packing + Repetition Matcher running on real pour schedule data"),
        ("📦", "Kitting",   "550 elements → auto-generated daily picking lists across 10 floors"),
        ("🔍", "Tracking",  "QR scan simulation updates panel health; auto-triggers shortage alerts"),
        ("📊", "BoQ Chart", "Interactive bar chart: baseline vs optimized per panel type with ₹ delta"),
        ("♻️", "Reuse",    "Panels marked REUSE/PARTIAL_REUSE/NEW_ORDER per pour, floor-by-floor"),
        ("🔔", "Alerts",   "Auto-flags SCRAP, HIGH_USAGE, NEEDS_REPAIR panels for replacement order"),
    ]
    for i, (icon, title, desc) in enumerate(features):
        x = Inches(0.3 + (i % 3) * 4.3)
        y = Inches(3.35 + (i // 3) * 1.7)
        add_rect(slide, x, y, Inches(4.0), Inches(1.5), RGBColor(0xFA,0xFA,0xFA))
        add_rect(slide, x, y, Inches(0.06), Inches(1.5), LT_RED)
        add_text_box(slide, f"{icon}  {title}", x + Inches(0.12), y + Inches(0.08),
                     Inches(3.8), Inches(0.4), size=13, bold=True, color=LT_NAVY)
        add_text_box(slide, desc, x + Inches(0.12), y + Inches(0.5),
                     Inches(3.8), Inches(0.9), size=11, color=LT_DGRAY)

    add_rect(slide, 0, Inches(7.1), W, Inches(0.4), LT_RED)

def slide_tech_stack(prs):
    slide = add_slide(prs)
    set_bg(slide, LT_WHITE)
    add_rect(slide, 0, 0, W, Inches(1.1), LT_RED)
    add_text_box(slide, "Technology Stack", Inches(0.5), Inches(0.2),
                 Inches(10), Inches(0.65), size=30, bold=True, color=LT_WHITE)

    tech = [
        ("React + Vite",       "Frontend SPA — Recharts, Framer Motion,\nReact Router, Axios"),
        ("Node.js + Express",  "REST API Gateway — Mongoose, dotenv,\nmulter file upload, CORS"),
        ("Python FastAPI",     "Optimization Engine — pandas, numpy,\nPydantic models, uvicorn"),
        ("MongoDB",           "Database — OptimizedKit, LiveInventory,\nFormworkCatalog collections"),
        ("2D Bin Packing",    "Core algorithm: greedy panel fill with\nremainder filler strips"),
        ("Repetition Matcher","Reuse scheduler: freed panel pool matched\nto upcoming requirements"),
    ]

    for i, (name, desc) in enumerate(tech):
        x = Inches(0.3 + (i % 3) * 4.3)
        y = Inches(1.3 + (i // 3) * 2.8)
        col = LT_NAVY if i % 2 == 0 else LT_RED
        add_rect(slide, x, y, Inches(4.0), Inches(2.5), RGBColor(0xF5,0xF8,0xFF))
        add_rect(slide, x, y, Inches(4.0), Inches(0.55), col)
        add_text_box(slide, name, x + Inches(0.1), y + Inches(0.07),
                     Inches(3.8), Inches(0.42), size=14, bold=True, color=LT_WHITE)
        add_text_box(slide, desc, x + Inches(0.1), y + Inches(0.65),
                     Inches(3.8), Inches(1.7), size=12, color=LT_DGRAY)

    add_rect(slide, 0, Inches(7.1), W, Inches(0.4), LT_NAVY)

def slide_closing(prs):
    slide = add_slide(prs)
    set_bg(slide, LT_NAVY)
    add_rect(slide, 0, 0, W, Inches(0.08), LT_RED)
    add_rect(slide, 0, Inches(7.42), W, Inches(0.08), LT_RED)

    add_text_box(slide, "FormworkIQ", Inches(1.0), Inches(1.5), Inches(11), Inches(1.4),
                 size=60, bold=True, color=LT_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "From spreadsheets to intelligence — one algorithm at a time.",
                 Inches(0.5), Inches(3.0), W - Inches(1.0), Inches(0.7),
                 size=20, italic=True, color=RGBColor(0xB0,0xC8,0xE8), align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(3.5), Inches(3.9), Inches(6.33), Inches(0.06), LT_RED)

    bullets = [
        "✅  Open-source stack — deployable on any cloud",
        "✅  Integrates with Revit BIM via JSON export",
        "✅  Scales to 1000+ element mega-structures",
        "✅  QR feedback loop ready for site deployment",
    ]
    for i, b in enumerate(bullets):
        add_text_box(slide, b, Inches(2.5), Inches(4.2 + i * 0.58), Inches(9), Inches(0.5),
                     size=15, color=LT_GRAY, align=PP_ALIGN.CENTER)

    add_text_box(slide, "L&T CreaTech 2026  ·  Problem Statement 4  ·  Formwork Kitting & BoQ Optimization",
                 Inches(0.3), Inches(7.1), W - Inches(0.6), Inches(0.4),
                 size=10, color=LT_LGRAY, align=PP_ALIGN.CENTER)

# ── Main ──────────────────────────────────────────────────────────

def generate_ppt(output_path="FormworkIQ_LT_CreaTech_2026.pptx"):
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print("Building slides...")
    slide_title(prs)
    print("  ✅ Slide 1: Title")
    slide_problem(prs)
    print("  ✅ Slide 2: Problem Statement")
    slide_solution(prs)
    print("  ✅ Slide 3: Our Solution")
    slide_architecture(prs)
    print("  ✅ Slide 4: System Architecture")
    slide_algorithm(prs)
    print("  ✅ Slide 5: Core Algorithm")
    slide_results(prs)
    print("  ✅ Slide 6: Results & Impact")
    slide_tech_stack(prs)
    print("  ✅ Slide 7: Technology Stack")
    slide_closing(prs)
    print("  ✅ Slide 8: Closing")

    prs.save(output_path)
    print(f"\n🎯 Presentation saved → {output_path}")
    print(f"   {prs.slides.__len__()} slides  |  16:9 widescreen  |  L&T brand theme")

if __name__ == "__main__":
    generate_ppt()
